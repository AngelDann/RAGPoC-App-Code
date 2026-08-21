from __future__ import annotations

from pathlib import Path
from typing import Any

OFFICE_SUFFIXES = {".docx", ".xlsx", ".xls", ".pptx"}


def extract_docx(path: Path) -> str:
    """Extract paragraphs and tables from a Word (.docx) document into Markdown."""
    import docx
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = docx.Document(path)
    lines: list[str] = []

    for element in doc.element.body:
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag
        if tag == "p":
            p = Paragraph(element, doc)
            text = p.text.strip()
            if not text:
                continue
            style_name = p.style.name if p.style else ""
            if style_name.startswith("Heading 1"):
                lines.append(f"# {text}\n")
            elif style_name.startswith("Heading 2"):
                lines.append(f"## {text}\n")
            elif style_name.startswith("Heading 3"):
                lines.append(f"### {text}\n")
            elif style_name.startswith("Heading 4"):
                lines.append(f"#### {text}\n")
            else:
                lines.append(f"{text}\n")
        elif tag == "tbl":
            table = Table(element, doc)
            table_md = _docx_table_to_markdown(table)
            if table_md:
                lines.append(table_md)

    return "\n".join(lines).strip()


def _docx_table_to_markdown(table: Any) -> str:
    rows = table.rows
    if not rows:
        return ""
    data: list[list[str]] = []
    for row in rows:
        row_vals = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        if any(row_vals):
            data.append(row_vals)
    if not data:
        return ""

    headers = data[0]
    num_cols = len(headers)
    output: list[str] = [
        "| " + " | ".join(headers) + " |",
        "| " + " | ".join(["---"] * num_cols) + " |",
    ]
    for row in data[1:]:
        padded = (row + [""] * num_cols)[:num_cols]
        output.append("| " + " | ".join(padded) + " |")

    return "\n".join(output) + "\n"


def extract_xlsx(path: Path, max_rows_per_sheet: int = 2000) -> str:
    """Extract worksheets and tables from Excel (.xlsx) into Markdown."""
    import openpyxl

    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    sheets_output: list[str] = []

    try:
        for sheetname in wb.sheetnames:
            ws = wb[sheetname]
            sheet_lines: list[str] = [f"## Hoja: {sheetname}\n"]
            rows_data: list[list[str]] = []
            row_count = 0

            for row in ws.iter_rows(values_only=True):
                row_count += 1
                if row_count > max_rows_per_sheet:
                    sheet_lines.append(f"\n*(Truncado a {max_rows_per_sheet} filas)*\n")
                    break
                row_strs = ["" if v is None else str(v).strip().replace("\n", " ") for v in row]
                if any(row_strs):
                    rows_data.append(row_strs)

            if not rows_data:
                continue

            headers = rows_data[0]
            num_cols = len(headers)
            table_md = [
                "| " + " | ".join(headers) + " |",
                "| " + " | ".join(["---"] * num_cols) + " |",
            ]
            for row in rows_data[1:]:
                padded = (row + [""] * num_cols)[:num_cols]
                table_md.append("| " + " | ".join(padded) + " |")

            sheet_lines.append("\n".join(table_md))
            sheets_output.append("\n".join(sheet_lines))
    finally:
        wb.close()

    return "\n\n".join(sheets_output).strip()


def extract_pptx(path: Path) -> str:
    """Extract slides, text boxes, tables and speaker notes from PowerPoint (.pptx)."""
    import pptx

    prs = pptx.Presentation(path)
    slides_output: list[str] = []

    for index, slide in enumerate(prs.slides, start=1):
        slide_lines: list[str] = []
        title = f"Diapositiva {index}"

        if slide.shapes.title and slide.shapes.title.text.strip():
            title = f"Diapositiva {index}: {slide.shapes.title.text.strip()}"

        slide_lines.append(f"## {title}\n")

        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        indent = "  " * paragraph.level
                        slide_lines.append(f"{indent}- {text}")
            elif shape.has_table:
                table_data: list[list[str]] = []
                for row in shape.table.rows:
                    row_vals = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    if any(row_vals):
                        table_data.append(row_vals)
                if table_data:
                    num_cols = len(table_data[0])
                    table_md = [
                        "| " + " | ".join(table_data[0]) + " |",
                        "| " + " | ".join(["---"] * num_cols) + " |",
                    ]
                    for row in table_data[1:]:
                        padded = (row + [""] * num_cols)[:num_cols]
                        table_md.append("| " + " | ".join(padded) + " |")
                    slide_lines.append("\n" + "\n".join(table_md) + "\n")

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_lines.append(f"\n> **Notas del orador:** {notes_text}")

        slides_output.append("\n".join(slide_lines))

    return "\n\n---\n\n".join(slides_output).strip()


def extract_office(path: Path) -> str:
    """Extract content from any supported office document."""
    suffix = path.suffix.lower()
    if suffix == ".docx":
        return extract_docx(path)
    if suffix in {".xlsx", ".xls"}:
        return extract_xlsx(path)
    if suffix == ".pptx":
        return extract_pptx(path)
    raise ValueError(f"Unsupported office file type: {suffix}")
