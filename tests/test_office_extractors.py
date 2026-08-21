from pathlib import Path

import docx
import openpyxl
import pptx

from ragpoc.extractors.office import extract_docx, extract_office, extract_pptx, extract_xlsx
from ragpoc.extractors.text import extract_text


def test_extract_docx(tmp_path: Path):
    doc_path = tmp_path / "sample.docx"
    doc = docx.Document()
    doc.add_heading("Título Principal", level=1)
    doc.add_paragraph("Este es un párrafo de introducción sobre RAG.")
    doc.add_heading("Sección de Datos", level=2)
    
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Producto"
    table.rows[0].cells[1].text = "Precio"
    table.rows[1].cells[0].text = "Servidor"
    table.rows[1].cells[1].text = "100 USD"
    
    doc.save(doc_path)

    extracted = extract_docx(doc_path)
    assert "# Título Principal" in extracted
    assert "Este es un párrafo de introducción sobre RAG." in extracted
    assert "## Sección de Datos" in extracted
    assert "| Producto | Precio |" in extracted
    assert "| Servidor | 100 USD |" in extracted

    # Test via extract_text
    assert extract_text(doc_path) == extracted


def test_extract_xlsx(tmp_path: Path):
    xlsx_path = tmp_path / "sample.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Ventas"
    ws.append(["Mes", "Ingresos", "Gastos"])
    ws.append(["Enero", 10000, 4000])
    ws.append(["Febrero", 12000, 4500])

    ws2 = wb.create_sheet(title="Usuarios")
    ws2.append(["ID", "Nombre"])
    ws2.append([1, "Alice"])

    wb.save(xlsx_path)

    extracted = extract_xlsx(xlsx_path)
    assert "## Hoja: Ventas" in extracted
    assert "| Mes | Ingresos | Gastos |" in extracted
    assert "| Enero | 10000 | 4000 |" in extracted
    assert "## Hoja: Usuarios" in extracted
    assert "| ID | Nombre |" in extracted
    assert "| 1 | Alice |" in extracted

    assert extract_text(xlsx_path) == extracted


def test_extract_pptx(tmp_path: Path):
    pptx_path = tmp_path / "sample.pptx"
    prs = pptx.Presentation()
    
    # Blank slide layout
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    txBox = slide.shapes.add_textbox(0, 0, 1000, 1000)
    tf = txBox.text_frame
    tf.text = "Arquitectura del Sistema"
    p = tf.add_paragraph()
    p.text = "Componente Backend en Django"

    # Add speaker note
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = "Explicar los modelos y servicios en detalle."

    prs.save(pptx_path)

    extracted = extract_pptx(pptx_path)
    assert "## Diapositiva 1" in extracted
    assert "Arquitectura del Sistema" in extracted
    assert "Componente Backend en Django" in extracted
    assert "Notas del orador:" in extracted
    assert "Explicar los modelos y servicios en detalle." in extracted

    assert extract_text(pptx_path) == extracted
